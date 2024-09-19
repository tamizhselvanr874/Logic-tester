import streamlit as st  
from pptx import Presentation  
from pptx.enum.shapes import MSO_SHAPE_TYPE  
import requests  
from io import BytesIO  
from PIL import Image  
import fitz  # PyMuPDF  
from docx import Document  
from docx.shared import Inches  
  
# URL of your Azure function endpoint  
azure_function_url = 'https://doc2pdf.azurewebsites.net/api/HttpTrigger1'  
  
# Function to convert PPT to PDF using Azure Function  
def ppt_to_pdf(ppt_file, pdf_file):  
    mime_type = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'  
    headers = {  
        "Content-Type": "application/octet-stream",  
        "Content-Type-Actual": mime_type  
    }  
    with open(ppt_file, 'rb') as file:  
        response = requests.post(azure_function_url, data=file.read(), headers=headers)  
        if response.status_code == 200:  
            with open(pdf_file, 'wb') as pdf_out:  
                pdf_out.write(response.content)  
            return True  
        else:  
            st.error(f"File conversion failed with status code: {response.status_code}")  
            st.error(f"Response: {response.text}")  
            return False  
  
# Logic to identify slides with visual elements  
def identify_visual_elements(ppt_file):  
    presentation = Presentation(ppt_file)  
    visual_slides = []  
    for slide_number, slide in enumerate(presentation.slides, start=1):  
        has_visual_elements = False  
        for shape in slide.shapes:  
            if shape.shape_type in {MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.TABLE, MSO_SHAPE_TYPE.CHART, MSO_SHAPE_TYPE.GROUP, MSO_SHAPE_TYPE.AUTO_SHAPE}:  
                has_visual_elements = True  
                break  
        if has_visual_elements:  
            visual_slides.append(slide_number)  
    return visual_slides  
  
# Function to capture images from the identified slides  
def capture_slide_images(pdf_file, slide_numbers):  
    doc = fitz.open(pdf_file)  
    images = []  
    for slide_number in slide_numbers:  
        page = doc[slide_number - 1]  
        pix = page.get_pixmap()  
        image = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)  
        buffer = BytesIO()  
        image.save(buffer, format="PNG")  
        images.append({"slide_number": slide_number, "image": buffer.getvalue()})  
    return images  
  
# Streamlit app  
def main():  
    st.title("PowerPoint Visual Elements Extractor (Images, Flowdiagrams, Tabels, ...)")  
      
    uploaded_file = st.file_uploader("Upload a PowerPoint file", type=['pptx'])  
      
    if uploaded_file is not None:  
        with open("uploaded_ppt.pptx", "wb") as f:  
            f.write(uploaded_file.read())  
          
        pdf_file = "converted.pdf"  
          
        if ppt_to_pdf("uploaded_ppt.pptx", pdf_file):  
            st.success("PPT to PDF conversion successful!")  
              
            visual_slide_numbers = identify_visual_elements("uploaded_ppt.pptx")  
              
            if visual_slide_numbers:  
                st.write(f"Slides with visual elements: {visual_slide_numbers}")  
                  
                images = capture_slide_images(pdf_file, visual_slide_numbers)  
                  
                doc = Document()  
                doc.add_heading('Slides with Visual Elements', level=1)  
                  
                for item in images:  
                    slide_number = item["slide_number"]  
                    image_data = item["image"]  
                      
                    st.image(image_data, caption=f"Slide {slide_number}", use_column_width=True)  
                      
                    doc.add_heading(f'Slide {slide_number}', level=2)  
                    image_stream = BytesIO(image_data)  
                    doc.add_picture(image_stream, width=Inches(6))  
                  
                doc_file = BytesIO()  
                doc.save(doc_file)  
                doc_file.seek(0)  
                  
                st.download_button(  
                    label="Download Word Document",  
                    data=doc_file,  
                    file_name="slides_with_visual_elements.docx",  
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"  
                )  
            else:  
                st.write("No slides with visual elements found.")  
        else:  
            st.error("PPT to PDF conversion failed.")  
  
if __name__ == "__main__":  
    main()  
